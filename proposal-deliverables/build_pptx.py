"""Build the 2-slide WBGT CrewSafe SG proposal deck (16:9)."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

AST = "/Users/zhongchengtiong/Documents/Obsidian Vault/AD_wbgt/proposal-deliverables/assets"
OUT = "/Users/zhongchengtiong/Documents/Obsidian Vault/AD_wbgt/proposal-deliverables/WBGT-CrewSafe-SG-Proposal.pptx"

NAVY  = RGBColor(0x1F,0x3B,0x57)
NAVY2 = RGBColor(0x2B,0x4A,0x6B)
BLUE  = RGBColor(0x2E,0x6F,0xB0)
ORANGE= RGBColor(0xE8,0x87,0x3A)
TEXT  = RGBColor(0x1A,0x27,0x33)
MUTED = RGBColor(0x6B,0x7A,0x88)
GREEN = RGBColor(0x2E,0x7D,0x46)
LIGHT = RGBColor(0xEA,0xF1,0xF8)
BAND  = RGBColor(0xF5,0xF8,0xFC)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LINE  = RGBColor(0xE1,0xE8,0xEF)
STEEL = RGBColor(0xC6,0xD6,0xE6)
SW, SH = 13.333, 7.5

prs = None
BLANK = None

def new_prs():
    global prs, BLANK
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def tb(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=1.0):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        p.line_spacing = line_sp
        for (t,sz,c,b) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.name = "Calibri"
    return box

def pic(s, path, x, y, w=None, h=None):
    iw, ih = Image.open(path).size
    if w and not h: h = w*ih/iw
    if h and not w: w = h*iw/ih
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h)), w, h

def pill(s, x, y, w, h, fill, text, color, size, bold=True):
    sp = rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame; tf.word_wrap=False
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Calibri"
    return sp

# =====================================================================
def build_slide1():
    s1 = slide()
    rect(s1, 0,0, SW, SH, WHITE)
    rect(s1, 0,0, SW, 1.28, NAVY)
    rect(s1, 0,1.28, SW, 0.06, ORANGE)
    rect(s1, 0.55,0.32, 0.62,0.62, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
    tb(s1, 0.55,0.30, 0.62,0.62, [[("W",26,WHITE,True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s1, 1.32,0.30, 8.6,0.55, [[("WBGT CrewSafe SG",30,WHITE,True)]])
    tb(s1, 1.34,0.86, 9.0,0.35, [[("Human-supervised heat-safety operations for small outdoor crews in Singapore",13,STEEL,False)]])
    pill(s1, 9.9,0.44, 2.9,0.44, NAVY2, "AD Project  ·  Sprint 0 Proposal  ·  Group 4", STEEL, 10.5)

    rect(s1, 0.55,1.62, 12.23,0.5, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    tb(s1, 0.75,1.62, 11.9,0.5,
       [[("Closed loop:   ",12,BLUE,True),
         ("weather signal  →  ML forecast  →  MOM policy  →  supervisor approval  →  worker action  →  audit evidence",12,TEXT,True)]],
       anchor=MSO_ANCHOR.MIDDLE)

    LX, LW = 0.55, 6.55
    tb(s1, LX,2.42, LW,0.3, [[("PROJECT DESCRIPTION",12,ORANGE,True)]])
    tb(s1, LX,2.74, LW,1.6,
       [[("WBGT CrewSafe SG turns Singapore's public NEA WBGT data into timely, MOM-aligned rest, "
          "hydration and task decisions for crews (landscaping, estate & campus maintenance, event setup) "
          "that refer to myENV instead of an on-site meter. NEA lightning-strike observations sit above the "
          "WBGT reading as an overriding stop-work hazard. A guarded AI agent drafts a plan from deterministic "
          "policy rules; the supervisor approves it; workers receive and acknowledge it on mobile — every step "
          "recorded end-to-end.",12,TEXT,False)]],
       line_sp=1.08)

    tb(s1, LX,4.44, LW,0.3, [[("VALUE TO STAKEHOLDERS",12,ORANGE,True)]])
    vals = [
        ("Outdoor workers", "clear one-tap instructions and acknowledgement — no wearables, no location tracking."),
        ("Site supervisors", "a live board that turns a heat signal into an approve-in-seconds decision, with proof of who received what."),
        ("Safety managers", "consistent MOM policy application plus an append-only audit trail as enforcement evidence."),
        ("Employer / operations", "closes the gap between seeing a reading and proving action — lower heat-illness and compliance risk."),
    ]
    vy = 4.78
    for head, body in vals:
        rect(s1, LX,vy+0.09, 0.12,0.12, ORANGE, shape=MSO_SHAPE.OVAL)
        tb(s1, LX+0.28,vy, LW-0.28,0.55,
           [[(head+"  —  ",11.5,TEXT,True),(body,11.5,MUTED,False)]], line_sp=1.03)
        vy += 0.60

    RX = 7.45
    rect(s1, RX,2.42, 5.33,4.82, WHITE, line=LINE, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    tb(s1, RX+0.28,2.58, 4.9,0.3, [[("HIGH-LEVEL USE CASE DIAGRAM",12,BLUE,True)]])
    _,pw,ph = pic(s1, f"{AST}/usecase.png", RX+0.55, 2.98, h=4.05)
    img = s1.shapes[-1]
    img.left = Inches(RX + (5.33-pw)/2)

# =====================================================================
def build_slide2():
    s2 = slide()
    rect(s2, 0,0, SW, SH, WHITE)
    rect(s2, 0,0, SW, 1.02, NAVY)
    rect(s2, 0,1.02, SW, 0.06, ORANGE)
    tb(s2, 0.55,0.20, 10.5,0.5, [[("Product Backlog  &  Prototype Screens",25,WHITE,True)]])
    tb(s2, 0.57,0.70, 10.0,0.3, [[("Proposed features, the technologies behind them, and Figma prototype screens",12.5,STEEL,False)]])
    pill(s2, 10.75,0.30, 2.05,0.42, NAVY2, "WBGT CrewSafe SG", STEEL, 11)

    tb(s2, 0.55,1.26, 8.2,0.3, [[("PRIORITISED PRODUCT BACKLOG   (extract — full 19-story backlog in the accompanying workbook)",11,ORANGE,True)]])
    rows = [
        ("Pri","Feature / user story","Key technologies","Pts","Spr"),
        ("Must","Authenticate & site-scoped access","React · React Native · Spring Security · JWT","5","1"),
        ("Must","Ingest WBGT / weather with freshness","Spring Boot · NEA data.gov.sg API · PostgreSQL","8","1"),
        ("Must","Lightning stop-work warning above WBGT","Spring Boot · NEA lightning API · React / React Native","5","1"),
        ("Must","Forecast 30 / 60-min WBGT","Python · FastAPI · scikit-learn / XGBoost","8","2"),
        ("Must","Evaluate deterministic heat policy","Java 21 rules engine · PostgreSQL config","8","2"),
        ("Must","Explainable agent draft plan","Tool-calling LLM · JSON-schema guarded tools","8","2"),
        ("Must","Dispatch & acknowledge worker action","React Native · Spring Boot · idempotency","8","2"),
        ("Must","Compliance & response-time dashboards","React · Chart.js · backend aggregation","8","3"),
        ("Must","Export audit timeline","Backend CSV / PDF · append-only audit","5","3"),
        ("Must","Safe degraded (NEA / ML / LLM) mode","Fallback adapters · persistence forecast","8","3"),
    ]
    tx, ty, tw = 0.55, 1.60, 8.15
    colw = [0.72, 3.05, 3.28, 0.5, 0.55]
    rh = 0.358
    scale = tw/sum(colw); colw=[c*scale for c in colw]
    for i, row in enumerate(rows):
        yy = ty + i*rh
        if i==0:
            rect(s2, tx, yy, tw, rh, NAVY)
        else:
            rect(s2, tx, yy, tw, rh, BAND if i%2==0 else WHITE)
        cx = tx
        for j, cell in enumerate(row):
            if i==0:
                tb(s2, cx+0.08, yy, colw[j]-0.12, rh, [[(cell,9.5,WHITE,True)]],
                   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER if j in (0,3,4) else PP_ALIGN.LEFT)
            elif j==0:
                pill(s2, cx+0.10, yy+0.055, colw[0]-0.20, rh-0.11, RGBColor(0xE7,0xF3,0xEC), cell, GREEN, 8.5)
            elif j in (3,4):
                tb(s2, cx, yy, colw[j], rh, [[(cell,9.5,TEXT,False)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            else:
                tb(s2, cx+0.08, yy, colw[j]-0.12, rh, [[(cell,9.3, TEXT if j==1 else MUTED, j==1)]], anchor=MSO_ANCHOR.MIDDLE)
            cx += colw[j]
        if i>0:
            rect(s2, tx, yy+rh-0.006, tw, 0.006, LINE)

    fy = ty + len(rows)*rh + 0.14
    rect(s2, tx, fy, tw, 0.74, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    tb(s2, tx+0.18, fy+0.10, tw-0.34, 0.6,
       [[("Cross-cutting:   ",10,BLUE,True),("Cloud — AWS Amplify Hosting (S3 + CloudFront) · ECS on Fargate · RDS for PostgreSQL · S3 · Secrets Manager",10,TEXT,False)],
        [("DevSecOps — ",10,BLUE,True),("GitHub Actions CI/CD · SAST + secret scan · dependency & container scanning · security-remediation evidence",10,TEXT,False)]],
       line_sp=1.05, sp_after=3)

    tb(s2, 8.85,1.26, 4.3,0.3, [[("PROTOTYPE SCREENS   (Figma)",11,ORANGE,True)]])
    gx0, colgap, gx1 = 8.85, 0.18, 12.98
    colw = (gx1-gx0-colgap)/2
    col1x, col2x = gx0, gx0+colw+colgap
    ph_h = 2.40; ph_w = ph_h*390/844
    wb_w = colw;  wb_h = wb_w*900/1440
    def gal(path, colx, y, iw, ih, cap):
        x = colx + (colw-iw)/2
        rect(s2, x-0.03, y-0.03, iw+0.06, ih+0.06, WHITE, line=LINE, lw=0.75)
        s2.shapes.add_picture(path, Inches(x), Inches(y), Inches(iw), Inches(ih))
        tb(s2, colx, y+ih+0.04, colw, 0.22, [[(cap,8.5,MUTED,True)]], align=PP_ALIGN.CENTER)
    ph_y, wb_y = 1.98, 4.84
    gal(f"{AST}/m1.png", col1x, ph_y, ph_w, ph_h, "Worker · readiness")
    gal(f"{AST}/m2.png", col2x, ph_y, ph_w, ph_h, "Worker · action")
    gal(f"{AST}/w1.png", col1x, wb_y, wb_w, wb_h, "Supervisor · live board")
    gal(f"{AST}/w2.png", col2x, wb_y, wb_w, wb_h, "Supervisor · approve plan")

# =====================================================================
if len(sys.argv) > 1 and sys.argv[1] == "preview":
    for idx, fn in [(1, build_slide1), (2, build_slide2)]:
        new_prs(); fn()
        p = OUT.replace(".pptx", f"__preview_s{idx}.pptx")
        prs.save(p); print("Saved preview:", p)
else:
    new_prs(); build_slide1(); build_slide2(); prs.save(OUT)
    print("Saved:", OUT)
